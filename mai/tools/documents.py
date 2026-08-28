"""Structured document-reading tools for PDF, DOCX, XLSX, and PPTX files."""
from __future__ import annotations

import os
from pathlib import Path
from typing import Any

from pydantic import BaseModel, ConfigDict, Field

from .registry import ToolRegistry


class DocumentReadInput(BaseModel):
    model_config = ConfigDict(extra="forbid")
    path: str = Field(min_length=1)
    max_chars: int = Field(default=50000, ge=1, le=500000)


def _resolve(path: str, cwd: str | Path | None) -> Path:
    candidate = Path(path).expanduser()
    if not candidate.is_absolute():
        candidate = Path(cwd or os.getcwd()) / candidate
    return candidate.resolve(strict=False)


def _trim(text: str, max_chars: int) -> tuple[str, bool]:
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars], True


def _read_pdf(path: Path) -> dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        pages.append({"page": index, "text": page.extract_text() or ""})
    return {"kind": "pdf", "pages": pages}


def _read_docx(path: Path) -> dict[str, Any]:
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs]
    tables = []
    for table in doc.tables:
        tables.append([[cell.text for cell in row.cells] for row in table.rows])
    return {"kind": "docx", "paragraphs": paragraphs, "tables": tables}


def _read_xlsx(path: Path) -> dict[str, Any]:
    from openpyxl import load_workbook

    workbook = load_workbook(filename=path, read_only=True, data_only=True)
    try:
        sheets = []
        for worksheet in workbook.worksheets:
            rows = [[cell.value for cell in row] for row in worksheet.iter_rows()]
            sheets.append({"title": worksheet.title, "rows": rows})
        return {"kind": "xlsx", "sheets": sheets}
    finally:
        workbook.close()


def _read_pptx(path: Path) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and text:
                texts.append(text)
        slides.append({"slide": index, "text": "\n".join(texts)})
    return {"kind": "pptx", "slides": slides}


def document_read(*, path: str, max_chars: int = 50000, cwd: str | Path | None = None) -> dict[str, Any]:
    target = _resolve(path, cwd)
    if not target.exists():
        raise FileNotFoundError(str(target))
    if not target.is_file():
        raise IsADirectoryError(str(target))

    suffix = target.suffix.lower()
    readers = {
        ".pdf": _read_pdf,
        ".docx": _read_docx,
        ".xlsx": _read_xlsx,
        ".pptx": _read_pptx,
    }
    reader = readers.get(suffix)
    if reader is None:
        raise ValueError(f"unsupported document extension: {suffix or '<none>'}")

    structured = reader(target)
    text_parts: list[str] = []
    details: dict[str, Any]
    kind = structured["kind"]
    if kind == "pdf":
        text_parts = [item["text"] for item in structured["pages"]]
        details = {"page_count": len(structured["pages"])}
    elif kind == "docx":
        text_parts = [*structured["paragraphs"], *["\t".join(row) for table in structured["tables"] for row in table]]
        details = {"paragraph_count": len(structured["paragraphs"]), "table_count": len(structured["tables"])}
    elif kind == "xlsx":
        text_parts = ["\t".join("" if value is None else str(value) for value in row) for sheet in structured["sheets"] for row in sheet["rows"]]
        details = {
            "sheets": [
                {"title": sheet["title"], "row_count": len(sheet["rows"])}
                for sheet in structured["sheets"]
            ]
        }
    elif kind == "pptx":
        text_parts = [item["text"] for item in structured["slides"]]
        details = {"slide_count": len(structured["slides"])}
    else:
        raise RuntimeError(f"unsupported structured document kind: {kind}")

    text, truncated = _trim("\n".join(text_parts), max_chars)
    return {
        "path": str(target),
        "extension": suffix,
        "text": text,
        "truncated": truncated,
        "details": details,
    }


def register_document_tools(
    registry: ToolRegistry,
    *,
    cwd: str | Path | None = None,
    timeout_seconds: float | None = 60,
) -> None:
    def handler(**kwargs: Any) -> Any:
        return document_read(cwd=cwd, **kwargs)

    registry.add(
        name="document_read",
        description=(
            "Read structured local document files. Supports PDF, DOCX, XLSX, and PPTX. "
            "Use this instead of file_read when the target is one of those binary document formats."
        ),
        input_model=DocumentReadInput,
        handler=handler,
        timeout_seconds=timeout_seconds,
        category="document",
    )
